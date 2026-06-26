from __future__ import annotations
import os
import io
import uuid
import copy
import json
import math
import threading
import traceback
import re
import time
from pathlib import Path
import requests

from flask import Flask, request, jsonify, send_file, render_template
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree
import anthropic

# ---------------------------------------------------------------------------
# App setup
# ---------------------------------------------------------------------------

app = Flask(__name__)

BASE_DIR = Path(__file__).parent
UPLOAD_FOLDER = BASE_DIR / "uploads"
OUTPUT_FOLDER = BASE_DIR / "outputs"
UPLOAD_FOLDER.mkdir(exist_ok=True)
OUTPUT_FOLDER.mkdir(exist_ok=True)

# In-memory job registry
# {job_id: {"status": ..., "message": ..., "progress": ..., "plan": ..., "pptx_path": ..., "output": ...}}
jobs: dict = {}
jobs_lock = threading.Lock()


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def clone_slide(prs: Presentation, source_index: int = 0):
    """Clone the slide at source_index and append a copy to the presentation."""
    source = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)

    for child in source.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    for rel in source.part.rels.values():
        if "image" in rel.reltype:
            try:
                new_slide.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                pass

    return new_slide


def replace_text_in_slide(slide, replacements: dict, ordered: dict = None):
    """
    Replace placeholder tokens in every text frame on a slide.

    replacements  – {old: new} for unique tokens
    ordered       – {old: [val1, val2, val3]} for tokens that appear
                    multiple times; replaced in document order (top→bottom)
    """
    order_counts = {k: 0 for k in (ordered or {})}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            if not full_text.strip():
                continue

            modified = full_text
            changed = False

            for placeholder, value in replacements.items():
                if placeholder in modified:
                    modified = modified.replace(
                        placeholder, str(value) if value is not None else ""
                    )
                    changed = True

            for placeholder, values in (ordered or {}).items():
                if placeholder in modified:
                    idx = order_counts[placeholder]
                    if idx < len(values):
                        modified = modified.replace(
                            placeholder,
                            str(values[idx]) if values[idx] is not None else "",
                        )
                        order_counts[placeholder] += 1
                        changed = True

            if changed and para.runs:
                para.runs[0].text = modified
                for run in para.runs[1:]:
                    run.text = ""


# ---------------------------------------------------------------------------
# Real landmark lookup — Google Maps (preferred) with OSM fallback
# ---------------------------------------------------------------------------

OSM_HEADERS = {"User-Agent": "Skyscale-OOH-Generator/1.0 (contact@skyscale.com)"}

def _haversine_km(lat1: float, lon1: float, lat2: float, lon2: float) -> float:
    R = 6371.0
    dlat = math.radians(lat2 - lat1)
    dlon = math.radians(lon2 - lon1)
    a = (math.sin(dlat / 2) ** 2
         + math.cos(math.radians(lat1)) * math.cos(math.radians(lat2))
         * math.sin(dlon / 2) ** 2)
    return R * 2 * math.asin(math.sqrt(a))


def _get_landmarks_google(location: str, city: str, n: int = 3):
    api_key = os.environ.get("GOOGLE_MAPS_API_KEY", "")
    if not api_key:
        return None
    try:
        geo = requests.get(
            "https://maps.googleapis.com/maps/api/geocode/json",
            params={"address": f"{location}, {city}", "key": api_key},
            timeout=8,
        ).json()
        if geo.get("status") != "OK":
            return None
        loc = geo["results"][0]["geometry"]["location"]
        lat, lng = loc["lat"], loc["lng"]

        places = requests.get(
            "https://maps.googleapis.com/maps/api/place/nearbysearch/json",
            params={
                "location": f"{lat},{lng}",
                "radius": 5000,
                "type": "point_of_interest",
                "key": api_key,
            },
            timeout=10,
        ).json()

        results = []
        seen: set = set()
        for p in places.get("results", []):
            name = p.get("name", "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            p_lat = p["geometry"]["location"]["lat"]
            p_lng = p["geometry"]["location"]["lng"]
            dist = _haversine_km(lat, lng, p_lat, p_lng)
            results.append((dist, name))

        results.sort(key=lambda x: x[0])
        filtered = [(d, name) for d, name in results if d <= 5.0][:n]
        return [
            f"{name} – {round(d, 1) if d >= 0.1 else 0.1}km"
            for d, name in filtered
        ] or None

    except Exception:
        return None


def _get_landmarks_osm(location: str, city: str, n: int = 3):
    try:
        geo_resp = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={"q": f"{location}, {city}", "format": "json", "limit": 1},
            headers=OSM_HEADERS,
            timeout=8,
        )
        geo_data = geo_resp.json()
        if not geo_data:
            return None
        lat = float(geo_data[0]["lat"])
        lon = float(geo_data[0]["lon"])
        time.sleep(1.1)

        overpass_query = f"""
[out:json][timeout:12];
(
  node["name"]["tourism"](around:5000,{lat},{lon});
  node["name"]["amenity"~"^(restaurant|cafe|hotel|bank|museum|theatre|cinema|hospital|university|library|historic)$"](around:5000,{lat},{lon});
  node["name"]["historic"](around:5000,{lat},{lon});
  node["name"]["shop"~"^(mall|department_store|supermarket)$"](around:5000,{lat},{lon});
);
out center 20;
"""
        elements = requests.post(
            "https://overpass-api.de/api/interpreter",
            data={"data": overpass_query},
            headers=OSM_HEADERS,
            timeout=15,
        ).json().get("elements", [])

        seen: set = set()
        ranked: list = []
        for el in elements:
            name = el.get("tags", {}).get("name", "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            el_lat = el.get("lat") or el.get("center", {}).get("lat", lat)
            el_lon = el.get("lon") or el.get("center", {}).get("lon", lon)
            ranked.append((_haversine_km(lat, lon, float(el_lat), float(el_lon)), name))

        ranked.sort(key=lambda x: x[0])
        results = []
        for dist, name in ranked:
            if dist > 5.0:
                break
            km = round(dist, 1) if dist >= 0.1 else 0.1
            results.append(f"{name} – {km}km")
            if len(results) == n:
                break

        return results if len(results) >= n else None

    except Exception:
        return None


def get_real_landmarks(location: str, city: str, n: int = 3):
    result = _get_landmarks_google(location, city, n)
    if result and len(result) >= n:
        return result
    return _get_landmarks_osm(location, city, n)


# ---------------------------------------------------------------------------
# Google Maps Static screenshot
# Adjust these constants to match your template layout (inches from top-left)
# ---------------------------------------------------------------------------

MAP_IMG_LEFT   = Inches(6.60)   # horizontal offset from slide left edge
MAP_IMG_TOP    = Inches(1.20)   # vertical offset from slide top edge
MAP_IMG_WIDTH  = Inches(3.10)   # image width
MAP_IMG_HEIGHT = Inches(2.10)   # image height


def get_map_image_bytes(location: str, city: str, zoom: int = 16) -> bytes | None:
    """
    Return a map image for the location.
    Uses Google Maps Static API when GOOGLE_MAPS_API_KEY is set,
    otherwise geocodes via Nominatim and fetches from staticmap.openstreetmap.de
    — no extra dependencies needed.
    """
    UA = "OOHProposalGenerator/1.0 (skyscalemedia.com)"
    api_key = os.environ.get("GOOGLE_MAPS_API_KEY", "")

    if api_key:
        try:
            geo = requests.get(
                "https://maps.googleapis.com/maps/api/geocode/json",
                params={"address": f"{location}, {city}", "key": api_key},
                timeout=8,
            ).json()
            if geo.get("status") != "OK" or not geo.get("results"):
                return None
            loc = geo["results"][0]["geometry"]["location"]
            lat, lng = loc["lat"], loc["lng"]
            resp = requests.get(
                "https://maps.googleapis.com/maps/api/staticmap",
                params={
                    "center": f"{lat},{lng}",
                    "zoom": zoom,
                    "size": "600x400",
                    "scale": 2,
                    "maptype": "roadmap",
                    "markers": f"color:red|size:mid|{lat},{lng}",
                    "key": api_key,
                },
                timeout=12,
            )
            ct = resp.headers.get("content-type", "")
            if resp.status_code == 200 and ct.startswith("image"):
                return resp.content
        except Exception:
            pass
        return None

    # ── Free fallback: Nominatim geocode → stitch OSM tiles with Pillow ──
    # Respect Nominatim's 1 req/sec rate limit
    time.sleep(1.1)
    try:
        geo = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={"q": f"{location}, {city}", "format": "json", "limit": 1},
            headers={"User-Agent": UA},
            timeout=8,
        ).json()
        if not geo:
            print(f"[MAP] Nominatim found nothing for {location!r}, {city!r}")
            return None
        lat, lng = float(geo[0]["lat"]), float(geo[0]["lon"])
        print(f"[MAP] geocoded {location!r} → lat={lat:.4f} lng={lng:.4f}")
    except Exception as e:
        print(f"[MAP] Nominatim failed for {location!r}: {e}")
        return None

    # Stitch a 3×3 grid of OSM tiles into one image
    try:
        from PIL import Image, ImageDraw

        z = min(zoom, 15)
        n = 2 ** z
        lat_rad = math.radians(lat)
        tx = int((lng + 180) / 360 * n)
        ty = int((1 - math.log(math.tan(lat_rad) + 1 / math.cos(lat_rad)) / math.pi) / 2 * n)

        TILE = 256
        canvas = Image.new("RGB", (TILE * 3, TILE * 3), (220, 220, 215))
        ok = 0
        for row in range(3):
            for col in range(3):
                url = f"https://tile.openstreetmap.org/{z}/{tx - 1 + col}/{ty - 1 + row}.png"
                try:
                    r = requests.get(url, headers={"User-Agent": UA}, timeout=8)
                    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
                        tile = Image.open(io.BytesIO(r.content)).convert("RGB")
                        canvas.paste(tile, (col * TILE, row * TILE))
                        ok += 1
                except Exception:
                    pass

        print(f"[MAP] stitched {ok}/9 tiles for {location!r}")

        # Red pin at centre
        draw = ImageDraw.Draw(canvas)
        cx, cy = canvas.width // 2, canvas.height // 2
        pr = 13
        draw.ellipse([cx - pr, cy - pr, cx + pr, cy + pr],
                     fill=(220, 30, 30), outline="white", width=3)

        buf = io.BytesIO()
        canvas.save(buf, format="PNG")
        return buf.getvalue()
    except Exception as e:
        print(f"[MAP] tile stitch failed for {location!r}: {e}")
        return None


def _cleanup_map_images(job_id: str):
    """Delete all map image files for a job."""
    with jobs_lock:
        map_paths = jobs.get(job_id, {}).get("map_paths", {})
    for path_str in map_paths.values():
        try:
            Path(path_str).unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# AI content generation
# ---------------------------------------------------------------------------

def generate_site_content(site: dict, client: anthropic.Anthropic) -> dict:
    site_name = site.get("Site Name", "")
    location  = site.get("Location", "")
    market    = site.get("Market", "")
    fmt       = site.get("Format", "")
    size      = site.get("Size", "")
    is_mobile = str(location).strip().lower() == "various"

    lookup_address = market if is_mobile else location
    real_landmarks: list | None = get_real_landmarks(lookup_address, market)

    if real_landmarks:
        landmark_instruction = (
            "Real nearby landmarks have already been sourced from a map service. "
            "For landmark_1/2/3 return exactly these strings unchanged:\n"
            + "\n".join(f"  {i+1}. {l}" for i, l in enumerate(real_landmarks))
        )
        # Use json.dumps to safely escape any special characters in landmark strings
        landmark_format = (
            f'"landmark_1": {json.dumps(real_landmarks[0])},\n'
            f'  "landmark_2": {json.dumps(real_landmarks[1])},\n'
            f'  "landmark_3": {json.dumps(real_landmarks[2])}'
        )
    else:
        landmark_instruction = (
            "Real map lookup was unavailable. Use your knowledge of this city to name "
            "3 specific, well-known nearby landmarks within 5km. "
            'Format each as "Landmark Name – 0.Xkm" (max 5km).'
        )
        landmark_format = (
            '"landmark_1": "Landmark Name – 0.Xkm",\n'
            '  "landmark_2": "Landmark Name – 0.Xkm",\n'
            '  "landmark_3": "Landmark Name – 0.Xkm"'
        )

    prompt = f"""You are writing punchy, professional copy for an OOH (Out-of-Home) advertising proposal.

Site details:
- Name: {site_name}
- Location / Address: {location}
- City / Market: {market}
- Format: {fmt}
- Size: {size}

Return ONLY valid JSON (no markdown fences, no extra text) with exactly these keys:

{{
  "tagline": "<4–7 word punchy advertising tagline for this site>",
  "location_desc": "<2–3 sentences describing where the site is and what surrounds it>",
  "visibility_desc": "<2–3 sentences about viewing angles, physical size, and sightlines>",
  "audience_desc": "<2–3 sentences about who passes by and approximate daily volume>",
  {landmark_format}
}}

{landmark_instruction}"""

    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=800,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    text = raw.strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        # Extract the outermost {...} block and retry
        m = re.search(r'\{[\s\S]*\}', text)
        if m:
            try:
                return json.loads(m.group())
            except json.JSONDecodeError:
                pass
        raise ValueError(f"AI returned invalid JSON: {text[:200]}")


# ---------------------------------------------------------------------------
# Helper: build replacement maps from a site plan dict
# ---------------------------------------------------------------------------

def _build_replacements(site: dict) -> tuple[dict, dict]:
    """Return (replacements, ordered) dicts for replace_text_in_slide."""
    lm = [site.get("landmark_1", ""), site.get("landmark_2", ""), site.get("landmark_3", "")]

    replacements = {
        # xyz-format template tokens
        "Site Name":       site.get("site_name", ""),
        "Headline":        site.get("tagline", ""),
        "Size: xyz":       f"Size: {site.get('size', '')}",
        "Format: xyz":     f"Format: {site.get('format', '')}",
        "Location: xyz":   f"Location: {site.get('location', '')}",
        "Frequency: xyz":  f"Frequency: {site.get('frequency', '')}",
        "Units: xyz":      f"Units: {site.get('units', '')}",
        "Traffic: xyz":    f"Traffic: {site.get('traffic', '')}",
        # {TOKEN} style
        "{SITE_NAME}":      site.get("site_name", ""),
        "{TAGLINE}":        site.get("tagline", ""),
        "{LOCATION_DESC}":  site.get("location_desc", ""),
        "{VISIBILITY_DESC}": site.get("visibility_desc", ""),
        "{AUDIENCE_DESC}":  site.get("audience_desc", ""),
        "{SIZE}":           site.get("size", ""),
        "{LOCATION}":       site.get("location", ""),
        "{UNITS}":          site.get("units", ""),
        "{FORMAT}":         site.get("format", ""),
        "{FREQUENCY}":      site.get("frequency", ""),
        "{TRAFFIC}":        site.get("traffic", ""),
        "{LANDMARK_1}":     lm[0],
        "{LANDMARK_2}":     lm[1],
        "{LANDMARK_3}":     lm[2],
        "{MARKET}":         site.get("market", ""),
    }

    ordered = {
        "Text": [
            site.get("location_desc", ""),
            site.get("visibility_desc", ""),
            site.get("audience_desc", ""),
        ],
        "Xyz - 0.5km":     lm,
        "Xyz – 0.5km": lm,
        "Xyz –0.5km":  lm,
        "Xyz -0.5km":      lm,
    }

    return replacements, ordered


# ---------------------------------------------------------------------------
# PPTX build worker (shared by both plan-based and legacy one-shot flows)
# ---------------------------------------------------------------------------

def build_pptx_from_plan(job_id: str, pptx_path: Path, plan: list):
    """Background job: build PPTX from a pre-computed plan list."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"] = status
            jobs[job_id]["message"] = message
            jobs[job_id]["progress"] = progress

    try:
        update("building", "Loading template…", 5)
        prs = Presentation(str(pptx_path))
        if not prs.slides:
            raise ValueError("The PowerPoint template has no slides.")

        template_slide  = prs.slides[0]
        template_layout = template_slide.slide_layout
        template_spTree = copy.deepcopy(template_slide.shapes._spTree)
        template_spTree_xml = etree.tostring(template_spTree, encoding="unicode")

        # Grab server-side map paths (not sent to client)
        with jobs_lock:
            map_paths = dict(jobs.get(job_id, {}).get("map_paths", {}))

        total = len(plan)

        for idx, site in enumerate(plan):
            pct = 5 + int((idx / total) * 90)
            update("building", f"Building slide {idx + 1}/{total}: {site.get('site_name', '')}…", pct)

            replacements, ordered = _build_replacements(site)

            if idx == 0:
                slide = prs.slides[0]
            else:
                slide = prs.slides.add_slide(template_layout)

                rId_map = {}
                for rel_id, rel in template_slide.part.rels.items():
                    if "image" in rel.reltype or "media" in rel.reltype:
                        try:
                            new_rId = slide.part.relate_to(rel.target_part, rel.reltype)
                            if new_rId != rel_id:
                                rId_map[rel_id] = new_rId
                        except Exception:
                            pass

                xml = template_spTree_xml
                for old_id, new_id in rId_map.items():
                    xml = xml.replace(f'r:embed="{old_id}"', f'r:embed="{new_id}"')
                    xml = xml.replace(f'r:link="{old_id}"',  f'r:link="{new_id}"')

                new_tree = slide.shapes._spTree
                for child in list(new_tree):
                    new_tree.remove(child)
                for child in etree.fromstring(xml):
                    new_tree.append(copy.deepcopy(child))

            replace_text_in_slide(slide, replacements, ordered)

            # Add Google Maps screenshot if available
            map_path_str = map_paths.get(idx)
            if map_path_str:
                map_file = Path(map_path_str)
                if map_file.exists():
                    try:
                        slide.shapes.add_picture(
                            str(map_file),
                            MAP_IMG_LEFT, MAP_IMG_TOP,
                            MAP_IMG_WIDTH, MAP_IMG_HEIGHT,
                        )
                        print(f"[PPTX] map added to slide {idx}")
                    except Exception as pic_err:
                        print(f"[WARN] add_picture failed for slide {idx}: {pic_err}")
                else:
                    print(f"[WARN] map file missing on disk for slide {idx}: {map_path_str}")
            else:
                print(f"[PPTX] no map path for slide {idx} — skipping")

        update("building", "Saving output file…", 96)
        output_filename = f"OOH_Proposal_{job_id[:8]}.pptx"
        output_path = OUTPUT_FOLDER / output_filename
        prs.save(str(output_path))

        with jobs_lock:
            jobs[job_id]["status"]   = "done"
            jobs[job_id]["message"]  = f"Done! {total} slide(s) generated."
            jobs[job_id]["progress"] = 100
            jobs[job_id]["output"]   = output_filename

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())

    finally:
        try:
            pptx_path.unlink(missing_ok=True)
        except Exception:
            pass
        _cleanup_map_images(job_id)


# ---------------------------------------------------------------------------
# Plan generation worker
# ---------------------------------------------------------------------------

def generate_plan_job(job_id: str, excel_path: Path):
    """Background job: read Excel + AI/landmarks, produce a content plan."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"]   = status
            jobs[job_id]["message"]  = message
            jobs[job_id]["progress"] = progress

    try:
        update("planning", "Reading Excel file…", 5)
        df = pd.read_excel(excel_path, engine="openpyxl")
        df.columns = [c.strip() for c in df.columns]

        if "Market" in df.columns:
            df["Market"] = df["Market"].ffill()

        if "Site Name" not in df.columns:
            raise ValueError("Excel file must have a 'Site Name' column.")

        df = df[df["Site Name"].notna() & (df["Site Name"].astype(str).str.strip() != "")]
        df = df.reset_index(drop=True)

        if df.empty:
            raise ValueError("No valid site rows found in the Excel file.")

        total = len(df)
        update("planning", f"Found {total} site(s). Connecting to AI…", 10)

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY environment variable is not set.")
        client = anthropic.Anthropic(api_key=api_key)

        plan = []
        for idx, row in df.iterrows():
            pct = 10 + int(((idx + 1) / total) * 85)
            site_name = str(row.get("Site Name", "")).strip()
            update("planning", f"Researching site {idx + 1}/{total}: {site_name}…", pct)

            location  = str(row.get("Location", "")).strip()
            market    = str(row.get("Market", "")).strip()
            is_mobile = location.lower() == "various"

            spot_dur = str(row.get("Spot Duration", "")).strip()
            sov_loop = str(row.get("SOV/Loop", "")).strip()
            if spot_dur.lower() in ("", "nan", "n/a", "na"):
                frequency = sov_loop
            else:
                frequency = f"{spot_dur} {sov_loop}".strip()

            raw_impacts = row.get("Impacts", "")
            try:
                traffic = f"{int(float(str(raw_impacts).replace(',', ''))):,}"
            except (ValueError, TypeError):
                traffic = str(raw_impacts).strip()

            try:
                ai = generate_site_content(row.to_dict(), client)
            except Exception as ai_err:
                print(f"[WARN] AI failed for site {site_name!r}: {ai_err}")
                ai = {
                    "tagline":         "",
                    "location_desc":   "",
                    "visibility_desc": "",
                    "audience_desc":   "",
                    "landmark_1":      "",
                    "landmark_2":      "",
                    "landmark_3":      "",
                }

            # Fetch map screenshot
            map_address = market if is_mobile else location
            map_zoom    = 11 if is_mobile else 16
            print(f"[MAP] fetching for site {idx} '{site_name}': address={map_address!r}, city={market!r}, zoom={map_zoom}")
            map_bytes   = get_map_image_bytes(map_address, market, zoom=map_zoom)
            has_map     = False
            if map_bytes:
                map_file = UPLOAD_FOLDER / f"{job_id}_map_{idx}.png"
                map_file.write_bytes(map_bytes)
                with jobs_lock:
                    jobs[job_id]["map_paths"][idx] = str(map_file)
                has_map = True
                print(f"[MAP] saved {len(map_bytes)} bytes for site {idx} '{site_name}'")
            else:
                print(f"[MAP] no map returned for site {idx} '{site_name}'")

            plan.append({
                "site_name":       site_name,
                "market":          market,
                "location":        location,
                "format":          str(row.get("Format", "")).strip(),
                "size":            str(row.get("Size", "")).strip(),
                "units":           str(row.get("Units/Faces", "")).strip(),
                "frequency":       frequency,
                "traffic":         traffic,
                "tagline":         ai.get("tagline", ""),
                "location_desc":   ai.get("location_desc", ""),
                "visibility_desc": ai.get("visibility_desc", ""),
                "audience_desc":   ai.get("audience_desc", ""),
                "landmark_1":      ai.get("landmark_1", ""),
                "landmark_2":      ai.get("landmark_2", ""),
                "landmark_3":      ai.get("landmark_3", ""),
                "has_map":         has_map,
            })

        with jobs_lock:
            jobs[job_id]["status"]   = "plan_ready"
            jobs[job_id]["message"]  = f"Plan ready — {total} site(s). Review and edit, then build."
            jobs[job_id]["progress"] = 100
            jobs[job_id]["plan"]     = plan

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())

    finally:
        try:
            excel_path.unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Legacy one-shot worker (kept for /api/generate backward compat)
# ---------------------------------------------------------------------------

def process_job(job_id: str, excel_path: Path, pptx_path: Path):
    """Upload + AI + build in one shot (legacy endpoint)."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"]   = status
            jobs[job_id]["message"]  = message
            jobs[job_id]["progress"] = progress

    try:
        update("processing", "Reading Excel file…", 5)
        df = pd.read_excel(excel_path, engine="openpyxl")
        df.columns = [c.strip() for c in df.columns]

        if "Market" in df.columns:
            df["Market"] = df["Market"].ffill()

        if "Site Name" not in df.columns:
            raise ValueError("Excel file must have a 'Site Name' column.")
        df = df[df["Site Name"].notna() & (df["Site Name"].astype(str).str.strip() != "")]
        df = df.reset_index(drop=True)

        if df.empty:
            raise ValueError("No valid site rows found in the Excel file.")

        total = len(df)
        update("processing", f"Found {total} site(s). Loading template…", 10)

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY environment variable is not set.")
        client = anthropic.Anthropic(api_key=api_key)

        plan = []
        for idx, row in df.iterrows():
            pct = 10 + int(((idx + 1) / total) * 80)
            site_name = str(row.get("Site Name", "")).strip()
            update("processing", f"Processing {idx + 1}/{total}: {site_name}…", pct)

            spot_dur = str(row.get("Spot Duration", "")).strip()
            sov_loop = str(row.get("SOV/Loop", "")).strip()
            frequency = sov_loop if spot_dur.lower() in ("", "nan", "n/a", "na") else f"{spot_dur} {sov_loop}".strip()

            raw_impacts = row.get("Impacts", "")
            try:
                traffic = f"{int(float(str(raw_impacts).replace(',', ''))):,}"
            except (ValueError, TypeError):
                traffic = str(raw_impacts).strip()

            ai = generate_site_content(row.to_dict(), client)

            plan.append({
                "site_name":       site_name,
                "market":          str(row.get("Market", "")).strip(),
                "location":        str(row.get("Location", "")).strip(),
                "format":          str(row.get("Format", "")).strip(),
                "size":            str(row.get("Size", "")).strip(),
                "units":           str(row.get("Units/Faces", "")).strip(),
                "frequency":       frequency,
                "traffic":         traffic,
                "tagline":         ai.get("tagline", ""),
                "location_desc":   ai.get("location_desc", ""),
                "visibility_desc": ai.get("visibility_desc", ""),
                "audience_desc":   ai.get("audience_desc", ""),
                "landmark_1":      ai.get("landmark_1", ""),
                "landmark_2":      ai.get("landmark_2", ""),
                "landmark_3":      ai.get("landmark_3", ""),
            })

        excel_path.unlink(missing_ok=True)
        build_pptx_from_plan(job_id, pptx_path, plan)

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())
        try:
            excel_path.unlink(missing_ok=True)
            pptx_path.unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return render_template("index.html")


# ── Pro flow ────────────────────────────────────────────────────────────────

@app.route("/api/plan", methods=["POST"])
def create_plan():
    """Step 1: Upload files, generate content plan (AI + landmarks)."""
    if "excel" not in request.files or "template" not in request.files:
        return jsonify({"error": "Both 'excel' and 'template' files are required."}), 400

    excel_file    = request.files["excel"]
    template_file = request.files["template"]

    if not excel_file.filename.endswith((".xlsx", ".xls")):
        return jsonify({"error": "Excel file must be .xlsx or .xls"}), 400
    if not template_file.filename.endswith(".pptx"):
        return jsonify({"error": "Template file must be .pptx"}), 400

    job_id     = uuid.uuid4().hex
    excel_path = UPLOAD_FOLDER / f"{job_id}_data.xlsx"
    pptx_path  = UPLOAD_FOLDER / f"{job_id}_template.pptx"
    excel_file.save(str(excel_path))
    template_file.save(str(pptx_path))

    with jobs_lock:
        jobs[job_id] = {
            "status":    "planning",
            "message":   "Starting…",
            "progress":  0,
            "plan":      None,
            "map_paths": {},        # {site_idx: str(path)} — server-side only
            "pptx_path": str(pptx_path),
            "output":    None,
        }

    threading.Thread(target=generate_plan_job, args=(job_id, excel_path), daemon=True).start()
    return jsonify({"job_id": job_id})


@app.route("/api/build", methods=["POST"])
def build():
    """Step 2: Submit (possibly edited) plan → build PPTX."""
    data = request.get_json(force=True, silent=True) or {}
    job_id = data.get("job_id")
    plan   = data.get("plan")

    if not job_id or not plan:
        return jsonify({"error": "job_id and plan are required"}), 400

    with jobs_lock:
        job = jobs.get(job_id)

    if not job:
        return jsonify({"error": "Job not found — session may have expired."}), 404

    pptx_path = Path(job.get("pptx_path", ""))
    if not pptx_path.exists():
        return jsonify({"error": "Template file not found. Please re-upload and start again."}), 404

    with jobs_lock:
        jobs[job_id]["status"]   = "building"
        jobs[job_id]["message"]  = "Starting build…"
        jobs[job_id]["progress"] = 0
        jobs[job_id]["output"]   = None

    threading.Thread(target=build_pptx_from_plan, args=(job_id, pptx_path, plan), daemon=True).start()
    return jsonify({"job_id": job_id})


# ── Shared polling + download ────────────────────────────────────────────────

@app.route("/api/status/<job_id>")
def status(job_id: str):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    # Only return fields safe for the frontend (omit pptx_path)
    return jsonify({
        "status":   job["status"],
        "message":  job["message"],
        "progress": job["progress"],
        "plan":     job.get("plan"),
        "output":   job.get("output"),
    })


@app.route("/api/download/<job_id>")
def download(job_id: str):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job or job["status"] != "done":
        return jsonify({"error": "File not ready"}), 404

    output_path = OUTPUT_FOLDER / job["output"]
    if not output_path.exists():
        return jsonify({"error": "Output file missing"}), 404

    response = send_file(
        str(output_path),
        as_attachment=True,
        download_name="OOH_Proposal.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    @response.call_on_close
    def _cleanup():
        try:
            output_path.unlink(missing_ok=True)
        except Exception:
            pass
        with jobs_lock:
            jobs.pop(job_id, None)

    return response


# ── Map image preview (served to the plan review UI) ────────────────────────

@app.route("/api/map/<job_id>/<int:site_idx>")
def serve_map(job_id: str, site_idx: int):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return "", 404
    map_path_str = job.get("map_paths", {}).get(site_idx)
    if not map_path_str:
        return "", 404
    p = Path(map_path_str)
    if not p.exists():
        return "", 404
    return send_file(str(p), mimetype="image/png")


# ── Legacy one-shot endpoint (backward compat) ───────────────────────────────

@app.route("/api/generate", methods=["POST"])
def generate():
    if "excel" not in request.files or "template" not in request.files:
        return jsonify({"error": "Both 'excel' and 'template' files are required."}), 400

    excel_file    = request.files["excel"]
    template_file = request.files["template"]

    if not excel_file.filename.endswith((".xlsx", ".xls")):
        return jsonify({"error": "Excel file must be .xlsx or .xls"}), 400
    if not template_file.filename.endswith(".pptx"):
        return jsonify({"error": "Template file must be .pptx"}), 400

    job_id     = uuid.uuid4().hex
    excel_path = UPLOAD_FOLDER / f"{job_id}_data.xlsx"
    pptx_path  = UPLOAD_FOLDER / f"{job_id}_template.pptx"
    excel_file.save(str(excel_path))
    template_file.save(str(pptx_path))

    with jobs_lock:
        jobs[job_id] = {
            "status":    "queued",
            "message":   "Queued…",
            "progress":  0,
            "plan":      None,
            "pptx_path": str(pptx_path),
            "output":    None,
        }

    threading.Thread(target=process_job, args=(job_id, excel_path, pptx_path), daemon=True).start()
    return jsonify({"job_id": job_id})


# ---------------------------------------------------------------------------
# CN Print Plan Filler
# ---------------------------------------------------------------------------

_CN_FIELDS = ['market','media','elements','format','platform','unit_type','kpis','buy_type','net_cpm','net_total']
_CN_COL = {'market':1,'media':2,'elements':3,'format':4,'platform':5,
            'unit_type':6,'kpis':7,'buy_type':8,'net_cpm':9,'net_total':10}
_CN_KWDS = {
    'market':    ['market','country','region','geo','territory'],
    'media':     ['site','publisher','media','outlet','supplied by','supplier','publication'],
    'elements':  ['package','description','placement','elements','product','brief'],
    'format':    ['placement name','format','ad format','format/specs','specs'],
    'platform':  ['platform','device','channel'],
    'unit_type': ['unit type','kpi type','metric','kpi'],
    'kpis':      ['kpi guarantee','units','quantity','impressions','views'],
    'buy_type':  ['cost method','buy type','pricing method','revenue type'],
    'net_cpm':   ['net cpm','net rate','cpm','rate'],
    'net_total': ['total usd','total net','net total','cost','total cost','investment','budget'],
}
_PLATFORM_ABBR = {'instagram':'IG','facebook':'FB','twitter':'X','x (twitter)':'X',
                  'linkedin':'LI','youtube':'YT','tiktok':'TT','snapchat':'SC'}

# Hardcoded column indices (0-indexed) for known CN rate card formats — from skill specification
_MARKET_MAPS = {
    'uk': {           # CN Traveller UK rate card
        'elements':  1,    # Col B(2): PLACEMENT
        'format':    2,    # Col C(3): FORMAT
        'unit_type': 7,    # Col H(8): KPI type
        'kpis':      12,   # Col M(13): KPI GUARANTEE
        'net_total': 17,   # Col R(18): TOTAL USD
    },
    'usa': {          # CN Traveller US rate card
        'media':     1,    # Col B(2): Site
        'elements':  4,    # Col E(5): Package
        'format':    5,    # Col F(6): Placement Name
        'platform':  6,    # Col G(7): Platform
        'unit_type': 8,    # Col I(9): Unit type
        'kpis':      9,    # Col J(10): Units
        'buy_type':  10,   # Col K(11): Cost method
        'net_cpm':   11,   # Col L(12): Rate
        'net_total': 12,   # Col M(13): Cost
    },
}


def _cn_match(header):
    h = str(header).lower().strip()
    if not h or h in ('nan','none'):
        return None
    for field, kws in _CN_KWDS.items():
        if h in kws:
            return field
    for field, kws in _CN_KWDS.items():
        for kw in kws:
            if kw in h or h in kw:
                return field
    return None


def _cn_find_header(rows):
    targets = {'site','placement','platform','cost','units','package','format','rate','cpm','media','kpi'}
    for i, row in enumerate(rows[:25]):
        vals = {str(v).lower().strip() for v in row if v is not None and str(v).strip()}
        if len(vals & targets) >= 2:
            return i
    return 0


def _cn_extract(file_bytes, market_label):
    import openpyxl as opx
    wb = opx.load_workbook(io.BytesIO(file_bytes), data_only=True)
    ws = wb.active
    raw = [[c.value for c in row] for row in ws.iter_rows()]
    if not raw:
        return []

    hdr = _cn_find_header(raw)
    headers = raw[hdr]
    data = raw[hdr + 1:]

    # Normalise market key for lookup
    mkt_key = market_label.lower().strip()
    if mkt_key in ('us', 'u.s.', 'u.s.a.', 'united states'):
        mkt_key = 'usa'
    hardcoded = _MARKET_MAPS.get(mkt_key, {})

    # Hardcoded columns first (reliable), then auto-detect the rest from headers
    field_col = {}
    used = set(hardcoded.values())
    for field in _CN_FIELDS:
        if field == 'market':
            continue
        if field in hardcoded:
            field_col[field] = hardcoded[field]
        else:
            for ci, h in enumerate(headers):
                if ci in used:
                    continue
                if _cn_match(h) == field:
                    field_col[field] = ci
                    used.add(ci)
                    break

    results, sparse_run, last_mkt = [], 0, ''
    for row in data:
        r = {}
        for field in _CN_FIELDS:
            ci = field_col.get(field)
            v = row[ci] if ci is not None and ci < len(row) else None
            s = str(v).strip() if v is not None else ''
            r[field] = '' if s in ('nan','None','N/A','n/a') else s

        r['market'] = r['market'] or market_label or last_mkt
        if r['market']:
            last_mkt = r['market']
        p = r.get('platform','')
        r['platform'] = _PLATFORM_ABBR.get(p.lower(), p)

        filled = sum(1 for f in _CN_FIELDS if f != 'market' and r[f])
        if filled < 3:
            sparse_run += 1
            if sparse_run >= 3 and results:
                break
            continue
        sparse_run = 0
        if any(len(r[f]) > 80 for f in _CN_FIELDS):
            continue
        results.append(r)

    return results


@app.route('/print-plan')
def print_plan_page():
    return render_template('print_plan.html')


@app.route('/fill-cn-plan', methods=['POST'])
def fill_cn_plan():
    import openpyxl as opx
    from openpyxl.cell.cell import MergedCell

    tpl = request.files.get('template')
    rcs = request.files.getlist('rate_cards')
    markets = request.form.getlist('markets')

    if not tpl:
        return 'Template file required', 400
    if not rcs or not any(f.filename for f in rcs):
        return 'At least one rate card required', 400

    try:
        wb = opx.load_workbook(io.BytesIO(tpl.read()))
        ws = wb.active

        # Clear values from row 6 onwards (preserve styles)
        for row in ws.iter_rows(min_row=6):
            for cell in row:
                if not isinstance(cell, MergedCell):
                    cell.value = None

        # Extract rows from each rate card
        all_rows = []
        for i, rc in enumerate(rcs):
            mkt = markets[i] if i < len(markets) else ''
            try:
                all_rows.extend(_cn_extract(rc.read(), mkt))
            except Exception as e:
                print(f'Rate card {rc.filename}: {e}')

        # Write into template starting at row 6
        last_mkt = ''
        for i, rd in enumerate(all_rows):
            r = 6 + i
            mkt = rd.get('market','')
            for field in _CN_FIELDS:
                col = _CN_COL[field]
                if field == 'market':
                    if mkt and mkt != last_mkt:
                        c = ws.cell(row=r, column=col)
                        if not isinstance(c, MergedCell):
                            c.value = mkt
                    continue
                val = rd.get(field,'')
                if not val:
                    continue
                c = ws.cell(row=r, column=col)
                if isinstance(c, MergedCell):
                    continue
                if field in ('net_cpm','net_total','kpis'):
                    try:
                        clean = re.sub(r'[^0-9.\-]','',val)
                        c.value = float(clean) if '.' in clean else int(clean)
                    except Exception:
                        c.value = val
                else:
                    c.value = val
            if mkt:
                last_mkt = mkt

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)
        base = tpl.filename.rsplit('.',1)[0] if tpl.filename else 'CN_Print_Plan'
        return send_file(out, as_attachment=True,
                         download_name=f'{base}_filled.xlsx',
                         mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

    except Exception as e:
        print(traceback.format_exc())
        return str(e), 500


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    port  = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_DEBUG", "0") == "1"
    app.run(host="0.0.0.0", port=port, debug=debug)
